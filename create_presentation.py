from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # --- CONSTANTS & THEME ---
    IBM_BLUE = RGBColor(15, 98, 254)      # #0F62FE
    DARK_BLUE = RGBColor(0, 29, 108)      # #001D6C
    COOL_GREY = RGBColor(244, 244, 244)   # #F4F4F4
    TEXT_BLACK = RGBColor(22, 22, 22)     # #161616
    WHITE = RGBColor(255, 255, 255)
    
    def apply_footer(slide, slide_num):
        """Adds standard footer to slide"""
        # Footer Line
        left = Inches(0.5)
        top = Inches(7.2)
        width = Inches(9)
        height = Inches(0.02)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COOL_GREY
        shape.line.fill.background()

        # Footer Text
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.25), Inches(9), Inches(0.25))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"IBM Project No. 14 | Syllabus & Curriculum Design Optimizer (SCDO) | Slide {slide_num}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(120, 120, 120)
        p.alignment = PP_ALIGN.RIGHT

    def apply_sidebar_accent(slide):
        """Adds a subtle vertical accent on the left"""
        left = Inches(0)
        top = Inches(0)
        width = Inches(0.15)
        height = Inches(7.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = IBM_BLUE
        shape.line.fill.background()

    def create_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        
        # Blue Header Shape
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2.5))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = DARK_BLUE
        header_shape.line.fill.background()
        
        # Accent Strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), Inches(10), Inches(0.1))
        strip.fill.solid()
        strip.fill.fore_color.rgb = IBM_BLUE
        strip.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(2))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.LEFT
        
        apply_footer(slide, 1)

    def add_content_slide(title_text, content_items, slide_num):
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        apply_sidebar_accent(slide)
        apply_footer(slide, slide_num)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.LEFT
        
        # Title Underline
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = IBM_BLUE
        line.line.fill.background()

        # Content Box
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        
        for item in content_items:
            p = tf.add_paragraph()
            p.space_after = Pt(12)
            
            if item.startswith("  -"):
                p.text = "  •  " + item.replace("  -", "").strip()
                p.level = 1
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(60, 60, 60)
            elif item.startswith("    *"):
                p.text = "      - " + item.replace("    *", "").strip()
                p.level = 2
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(90, 90, 90)
            else:
                p.text = item.strip()
                p.level = 0
                p.font.size = Pt(22)
                p.font.color.rgb = TEXT_BLACK
                # Bold logic for "Label:" patterns
                if ":" in item and len(item.split(":")[0]) < 35:
                     p.font.bold = True

    # --- SLIDES GENERATION ---

    # 1. Title
    create_title_slide(
        "Syllabus & Curriculum Design Optimizer",
        "Agentic AI-Powered System for Academic Compliance & Optimization\nIBM Project No. 14"
    )

    # 2. Problem Description
    add_content_slide("Problem Description", [
        "Manual Inefficiency:",
        "  - Curriculum design takes 40-60 faculty hours per course.",
        "Compliance Gaps:",
        "  - Difficulty aligning with dynamic NEP 2020, NBA, and NAAC standards.",
        "Static Content:",
        "  - Syllabi often lag behind industry trends (2-3 year gap).",
        "  - Lack of data-driven insights for improvements."
    ], 2)

    # 3. Objectives
    add_content_slide("Project Objectives", [
        "Automate:",
        "  - Extraction and structuring of legacy syllabus data (PDF/DOCX).",
        "Optimize:",
        "  - Enhance content using Bloom's Taxonomy and industry trends.",
        "Validate:",
        "  - Ensure 100% compliance with accreditation bodies via RAG.",
        "Standardize:",
        "  - Generate uniform, high-quality PDF/DOCX reports."
    ], 3)

    # 4. System Architecture
    add_content_slide("System Architecture Overview", [
        "Client-Server Model:",
        "  - React Frontend (Interactive Gap Analysis) ↔ FastAPI Backend.",
        "Microservices Architecture:",
        "  - Modular Parsers, Gap Analyzers, and Generation Engines.",
        "Data Flow:",
        "  - User Upload -> Parsing -> Vector Context (RAG) -> LLM Processing -> Structured Output."
    ], 4)

    # 5. Parsing Module
    add_content_slide("Module: Intelligent Parsing", [
        "Hybrid Approach:",
        "  - Uses `pdfplumber` for layout preservation and `PyPDF2` for raw text.",
        "Regex Pattern Matching:",
        "  - Detects Course Codes (`[A-Z]{2,}\d+`), Credits (`L-T-P`), and Units.",
        "Structure Recovery:",
        "  - Reconstructs hierarchy (Course -> Units -> Topics) from unstructured text."
    ], 5)

    # 6. RAG Engine
    add_content_slide("Module: RAG Analysis Engine", [
        "Retrieval Augmented Generation (RAG):",
        "  - 'Grounds' AI analysis in official manuals (NBA/NAAC) to prevent hallucinations.",
        "Vector Store:",
        "  - ChromaDB stores embeddings of accreditation guidelines and standard syllabi.",
        "Query Process:",
        "  - `RAGEngine` retrieves relevant 'Program Outcomes' and 'Rules' to validate input."
    ], 6)

    # 7. Content Optimization
    add_content_slide("Module: Content Optimization", [
        "Bloom's Taxonomy Mapper:",
        "  - Classifies outcomes (Remembering -> Creating) to ensure cognitive depth.",
        "Sequence Optimization:",
        "  - LLM re-orders units based on prerequisite dependency logic.",
        "Redundancy Check:",
        "  - Semantic similarity detection to remove overlapping topics across units."
    ], 7)

    # 8. Outcome Mapping
    add_content_slide("Module: Outcome Mapping", [
        "CO-PO Matrix Generation:",
        "  - Auto-generates correlation matrix (High/Med/Low) for accreditation.",
        "Logic:",
        "  - Calculates semantic cosine similarity between Course Outcomes (CO) and Program Outcomes (PO).",
        "Justification:",
        "  - Generates text explaining *why* a mapping exists."
    ], 8)

    # 9. Prompt Engineering
    add_content_slide("Methodology: Prompt Engineering", [
        "Role-Prompting:",
        "  - 'You are a curriculum expert with 20+ years experience...'",
        "Few-Shot Learning:",
        "  - Providing examples of 'Poor' vs 'Excellent' outcomes in system prompts.",
        "Structured Output Enforcement:",
        "  - Enforcing JSON or specific markdown formats to ensure parsing reliability."
    ], 9)

    # 10. Tech Stack
    add_content_slide("Technology Stack", [
        "Frontend:",
        "  - React.js, Tailwind CSS (Responsive Dashboard)",
        "Backend:",
        "  - FastAPI, Pydantic, Uvicorn (Async High Performance)",
        "AI/LLM:",
        "  - IBM Granite (via WatsonX.ai), LangChain (Orchestration)",
        "Data & Infrastructure:",
        "  - ChromaDB (Vectors), IBM Cloud Object Storage"
    ], 10)

    # 11. Results 1
    add_content_slide("Results: Content Optimization", [
        "Input (Weak Outcome):",
        "  - 'Understand Java'",
        "AI Output (Optimized):",
        "  - 'Design and implement object-oriented applications using Java collections and exception handling.'",
        "Improvement:",
        "  - Shifted from 'Understand' (Low Bloom) to 'Create' (High Bloom).",
        "  - Added specific technical criteria (Collections, Exception Handling)."
    ], 11)

    # 12. Results 2
    add_content_slide("Results: Gap Analysis", [
        "Detected Gap:",
        "  - Missing 'Professional Ethics' module required by NBA for engineering courses.",
        "System Action:",
        "  - Suggests insertion of 'Unit 6: Professional Ethics and Cyber Law'.",
        "  - Provides standard topics and learning resources for the new unit."
    ], 12)

    # 13. Impact
    add_content_slide("Impact & Novelty", [
        "Novelty:",
        "  - First Agentic RAG system specifically tuned for Indian Academic Standards (NEP/NBA).",
        "Impact:",
        "  - Reduces syllabus design time by ~80%.",
        "  - Increases compliance accuracy to near 100%.",
        "  - Standardizes quality across departments."
    ], 13)

    # 14. UI
    add_content_slide("User Interface", [
        "Interactive Dashboard:",
        "  - Drag-and-drop file upload.",
        "  - Real-time editing of parsed content.",
        "  - Visual Gap Analysis reports with color-coded alerts.",
        "  - (Live Demo / Screenshots available in project folder)"
    ], 14)

    # 15. Future Scope
    add_content_slide("Future Scope", [
        "LMS Integration:",
        "  - Push approved syllabi directly to Moodle/Canvas APIs.",
        "Real-Time Data:",
        "  - Live API keys to fetch current job descriptions (LinkedIn/Indeed) for dynamic skill mapping.",
        "Multilingual Support:",
        "  - Support for regional language syllabi processing."
    ], 15)

    # 16. Conclusion
    create_title_slide("Conclusion", "SCDO successfully bridges the gap between static academic planning and dynamic industry needs.")

    output_path = "Syllabus_Optimizer_Presentation_Enhanced.pptx"
    try:
        prs.save(output_path)
        print(f"Presentation saved to {output_path}")
    except PermissionError:
        print(f"Error: Could not save to {output_path}. Is the file open?")
        # Fallback
        prs.save("Syllabus_Optimizer_Presentation_Final.pptx")
        print("Saved to Syllabus_Optimizer_Presentation_Final.pptx instead.")

if __name__ == "__main__":
    create_presentation()
