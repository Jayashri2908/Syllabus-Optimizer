from pptx import Presentation
import os

def analyze_ppt(file_path):
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return

    prs = Presentation(file_path)
    
    print(f"Analyzing: {file_path}")
    print("-" * 30)
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    text_runs.append(paragraph.text)
        
        # print non-empty lines
        for line in text_runs:
            if line.strip():
                print(line.strip())
if __name__ == "__main__":
    analyze_ppt("Syllabus_Optimizer_Presentation_Enhanced.pptx")
